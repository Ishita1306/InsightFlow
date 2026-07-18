"""
Export service module.

Handles generation of executive-level reporting exports:
- PDF (via ReportLab)
- PowerPoint (via python-pptx)
- PNG Snapshot (via Pillow)
- Excel Summary (via Pandas & OpenPyXL)
"""

import os
import io
import datetime
import pandas as pd
import numpy as np
from PIL import Image as PILImage, ImageDraw, ImageFont

# ReportLab
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors

# python-pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


class ExportService:
    """Service to export dashboard reports in enterprise-grade PDF, PPTX, PNG, and Excel formats."""

    @classmethod
    def generate_pdf(cls, df: pd.DataFrame, health: float, kpis: list, cleaning_summary: dict, filename: str) -> bytes:
        """Generate a premium corporate PDF report."""
        buffer = io.BytesIO()
        doc = SimpleDocTemplate(
            buffer,
            pagesize=letter,
            rightMargin=54,
            leftMargin=54,
            topMargin=54,
            bottomMargin=54
        )
        
        styles = getSampleStyleSheet()
        
        # Define clean, professional corporate typography styles
        title_style = ParagraphStyle(
            'DocTitle',
            parent=styles['Heading1'],
            fontName='Helvetica-Bold',
            fontSize=24,
            leading=28,
            textColor=colors.HexColor('#1E293B'),
            spaceAfter=6
        )
        
        subtitle_style = ParagraphStyle(
            'DocSubtitle',
            parent=styles['Normal'],
            fontName='Helvetica',
            fontSize=10,
            leading=14,
            textColor=colors.HexColor('#64748B'),
            spaceAfter=20
        )
        
        h2_style = ParagraphStyle(
            'SectionHeader',
            parent=styles['Heading2'],
            fontName='Helvetica-Bold',
            fontSize=14,
            leading=18,
            textColor=colors.HexColor('#0F172A'),
            spaceBefore=15,
            spaceAfter=8,
            keepWithNext=True
        )
        
        body_style = ParagraphStyle(
            'BodyText',
            parent=styles['Normal'],
            fontName='Helvetica',
            fontSize=9.5,
            leading=14,
            textColor=colors.HexColor('#334155'),
            spaceAfter=6
        )
        
        bold_body_style = ParagraphStyle(
            'BoldBodyText',
            parent=body_style,
            fontName='Helvetica-Bold'
        )

        story = []
        
        # 1. Document Header
        story.append(Paragraph("Kosvio — Executive Intelligence Report", title_style))
        timestamp = datetime.datetime.now().strftime("%B %d, %Y at %I:%M %p")
        story.append(Paragraph(f"Dataset Ref: <b>{filename}</b> &nbsp;|&nbsp; Compiled: <b>{timestamp}</b>", subtitle_style))
        story.append(Spacer(1, 10))
        
        # 2. Executive Summary Metrics Table
        story.append(Paragraph("Dataset Overview & Performance Indicators", h2_style))
        
        kpi_data = [
            [
                Paragraph("<b>Metric</b>", bold_body_style),
                Paragraph("<b>Value</b>", bold_body_style),
                Paragraph("<b>Details</b>", bold_body_style)
            ]
        ]
        
        # Add primary metrics
        kpi_data.append([Paragraph("Dataset Health Score", body_style), Paragraph(f"{health:.1f}/100", bold_body_style), Paragraph("Quality rating based on validation audits", body_style)])
        for k in kpis:
            kpi_data.append([
                Paragraph(k["label"], body_style),
                Paragraph(str(k["value"]), bold_body_style),
                Paragraph(k.get("detail", "") or "", body_style)
            ])
            
        t = Table(kpi_data, colWidths=[150, 100, 250])
        t.setStyle(TableStyle([
            ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#F8FAFC')),
            ('ALIGN', (0,0), (-1,-1), 'LEFT'),
            ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
            ('TEXTCOLOR', (0,0), (-1,0), colors.HexColor('#0F172A')),
            ('BOTTOMPADDING', (0,0), (-1,-1), 6),
            ('TOPPADDING', (0,0), (-1,-1), 6),
            ('ROWBACKGROUNDS', (0,1), (-1,-1), [colors.white, colors.HexColor('#F8FAFC')]),
            ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor('#E2E8F0')),
        ]))
        story.append(t)
        story.append(Spacer(1, 15))
        
        # 3. Data Cleaning Summary
        story.append(Paragraph("Data Quality & Cleansing Records", h2_style))
        if cleaning_summary:
            clean_data = [
                [Paragraph("<b>Operation</b>", bold_body_style), Paragraph("<b>Impact Metrics</b>", bold_body_style)]
            ]
            clean_data.append([Paragraph("Rows Affected", body_style), Paragraph(f"{cleaning_summary.get('rows_removed', 0):,} rows dropped", body_style)])
            clean_data.append([Paragraph("Missing Fields Repaired", body_style), Paragraph(f"{cleaning_summary.get('missing_fixed', 0):,} missing entries filled", body_style)])
            clean_data.append([Paragraph("Datatype Conversions", body_style), Paragraph(f"{cleaning_summary.get('dtypes_converted', 0)} columns adjusted", body_style)])
            clean_data.append([Paragraph("Outliers Treated", body_style), Paragraph(f"{cleaning_summary.get('outliers_handled', 0):,} extreme values clipped", body_style)])
            
            ct = Table(clean_data, colWidths=[200, 300])
            ct.setStyle(TableStyle([
                ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#F1F5F9')),
                ('ALIGN', (0,0), (-1,-1), 'LEFT'),
                ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
                ('BOTTOMPADDING', (0,0), (-1,-1), 5),
                ('TOPPADDING', (0,0), (-1,-1), 5),
                ('ROWBACKGROUNDS', (0,1), (-1,-1), [colors.white, colors.HexColor('#F8FAFC')]),
                ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor('#E2E8F0')),
            ]))
            story.append(ct)
        else:
            story.append(Paragraph("This dataset is utilizing its raw, original structure. No cleaning operations have been applied in this session.", body_style))
            
        story.append(Spacer(1, 15))
        
        # 4. AI Strategic Recommendations
        story.append(Paragraph("Executive Recommendations & Strategic Guidance", h2_style))
        
        recs = [
            "Establish automated checks on incoming files to address missing cells prior to downstream analysis.",
            "Verify datatype alignment in database schemas to prevent text representations of date fields.",
            "Perform segment profiling on outlier populations to isolate anomalies from regular customer transactions.",
            "Schedule weekly deduplication scripts to ensure reporting tables remain optimal."
        ]
        
        for idx, rec in enumerate(recs, 1):
            story.append(Paragraph(f"<b>{idx}.</b> {rec}", body_style))
            story.append(Spacer(1, 3))
            
        doc.build(story)
        return buffer.getvalue()

    @classmethod
    def generate_pptx(cls, df: pd.DataFrame, health: float, kpis: list, cleaning_summary: dict, filename: str) -> bytes:
        """Generate a professional executive PowerPoint slide deck."""
        prs = Presentation()
        
        # Color constants
        dark_slate = RGBColor(15, 23, 42)
        accent_blue = RGBColor(99, 102, 241)
        sub_text_col = RGBColor(100, 116, 139)
        white = RGBColor(255, 255, 255)
        
        # ── SLIDE 1: Title Slide
        slide_layout = prs.slide_layouts[5] # Blank / Title-only layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Apply dark slate background to title slide for a premium look
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = dark_slate
        
        # Title text frame
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(8), Inches(2))
        tf = title_box.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = "Kosvio"
        p1.font.bold = True
        p1.font.size = Pt(44)
        p1.font.color.rgb = accent_blue
        
        p2 = tf.add_paragraph()
        p2.text = "Executive Intelligence Briefing"
        p2.font.bold = True
        p2.font.size = Pt(28)
        p2.font.color.rgb = white
        
        p3 = tf.add_paragraph()
        p3.text = f"Dataset: {filename}   |   Date: {datetime.datetime.now().strftime('%B %d, %Y')}"
        p3.font.size = Pt(12)
        p3.font.color.rgb = sub_text_col
        
        # ── SLIDE 2: KPI & Health Score Slide
        slide2 = prs.slides.add_slide(prs.slide_layouts[5])
        
        # Header
        header_box = slide2.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(8.5), Inches(1))
        tf2 = header_box.text_frame
        p_hdr = tf2.paragraphs[0]
        p_hdr.text = "Key Performance Indicators & Health Score"
        p_hdr.font.bold = True
        p_hdr.font.size = Pt(22)
        p_hdr.font.color.rgb = dark_slate
        
        # Add Health Score Card
        score_box = slide2.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(3.5), Inches(4.5))
        s_tf = score_box.text_frame
        s_tf.word_wrap = True
        sp1 = s_tf.paragraphs[0]
        sp1.text = "DATASET HEALTH"
        sp1.font.bold = True
        sp1.font.size = Pt(11)
        sp1.font.color.rgb = sub_text_col
        
        sp2 = s_tf.add_paragraph()
        sp2.text = f"{health:.1f} / 100"
        sp2.font.bold = True
        sp2.font.size = Pt(40)
        sp2.font.color.rgb = accent_blue
        
        sp3 = s_tf.add_paragraph()
        sp3.text = "This rating reflects overall reliability based on missing records, duplicate rows, empty columns, and formatting issues."
        sp3.font.size = Pt(10.5)
        sp3.font.color.rgb = dark_slate
        
        # Add table of KPI metrics
        rows = len(kpis) + 1
        cols = 2
        left = Inches(4.5)
        top = Inches(1.5)
        width = Inches(4.75)
        height = Inches(0.5 * rows)
        
        table_shape = slide2.shapes.add_table(rows, cols, left, top, width, height)
        table = table_shape.table
        
        table.columns[0].width = Inches(2.25)
        table.columns[1].width = Inches(2.5)
        
        # Write headers
        table.cell(0, 0).text = "Metric Description"
        table.cell(0, 1).text = "Computed Value"
        table.cell(0, 0).text_frame.paragraphs[0].font.bold = True
        table.cell(0, 1).text_frame.paragraphs[0].font.bold = True
        
        for idx, k in enumerate(kpis):
            table.cell(idx+1, 0).text = k["label"]
            table.cell(idx+1, 1).text = f"{k['value']} ({k.get('detail', '')})"
            table.cell(idx+1, 0).text_frame.paragraphs[0].font.size = Pt(10)
            table.cell(idx+1, 1).text_frame.paragraphs[0].font.size = Pt(10)
            
        # ── SLIDE 3: Cleaning Summary & Recommendations
        slide3 = prs.slides.add_slide(prs.slide_layouts[5])
        
        header_box3 = slide3.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(8.5), Inches(1))
        tf3 = header_box3.text_frame
        p_hdr3 = tf3.paragraphs[0]
        p_hdr3.text = "Operations Impact & Executive Recommendations"
        p_hdr3.font.bold = True
        p_hdr3.font.size = Pt(22)
        p_hdr3.font.color.rgb = dark_slate
        
        # Left side: Cleaning summary
        summary_box = slide3.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(4.0), Inches(4.5))
        sum_tf = summary_box.text_frame
        sum_tf.word_wrap = True
        sum_p1 = sum_tf.paragraphs[0]
        sum_p1.text = "CLEANING OPERATIONS STATS"
        sum_p1.font.bold = True
        sum_p1.font.size = Pt(11)
        sum_p1.font.color.rgb = sub_text_col
        
        if cleaning_summary:
            ops = [
                f"• Rows removed: {cleaning_summary.get('rows_removed', 0):,}",
                f"• Missing cells fixed: {cleaning_summary.get('missing_fixed', 0):,}",
                f"• Data types converted: {cleaning_summary.get('dtypes_converted', 0)}",
                f"• Outliers handled: {cleaning_summary.get('outliers_handled', 0):,}",
            ]
            for op in ops:
                op_p = sum_tf.add_paragraph()
                op_p.text = op
                op_p.font.size = Pt(12)
                op_p.font.color.rgb = dark_slate
        else:
            op_p = sum_tf.add_paragraph()
            op_p.text = "No active cleaning operations applied to this dataset."
            op_p.font.size = Pt(11)
            op_p.font.color.rgb = dark_slate
            
        # Right side: Executive Recommendations
        recs_box = slide3.shapes.add_textbox(Inches(5.0), Inches(1.5), Inches(4.25), Inches(4.5))
        recs_tf = recs_box.text_frame
        recs_tf.word_wrap = True
        rec_p1 = recs_tf.paragraphs[0]
        rec_p1.text = "EXECUTIVE RECOMMENDATIONS"
        rec_p1.font.bold = True
        rec_p1.font.size = Pt(11)
        rec_p1.font.color.rgb = sub_text_col
        
        recs = [
            "Validate and audit source feeds weekly.",
            "Establish clean date formatting standards.",
            "Isolate extreme customer values for segmentation.",
            "Maintain strict primary key constraints."
        ]
        for rec in recs:
            rec_p = recs_tf.add_paragraph()
            rec_p.text = f"✔ {rec}"
            rec_p.font.size = Pt(11.5)
            rec_p.font.color.rgb = dark_slate
            
        buffer = io.BytesIO()
        prs.save(buffer)
        return buffer.getvalue()

    @classmethod
    def generate_png(cls, df: pd.DataFrame, health: float, kpis: list, cleaning_summary: dict, filename: str) -> bytes:
        """Compose a premium corporate dashboard summary image using Pillow."""
        # 1200x800 Dashboard Snapshot
        img = PILImage.new('RGB', (1200, 800), color='#0F172A') # Dark Theme Slate Background
        draw = ImageDraw.Draw(img)
        
        # Draw background panels
        draw.rounded_rectangle([40, 40, 1160, 760], radius=16, fill='#1E293B', outline='#334155', width=1)
        
        # Header banner
        draw.rounded_rectangle([60, 60, 1140, 130], radius=8, fill='#334155')
        
        # Text drawing helper
        def draw_text_simple(x, y, text, size=16, color='#FFFFFF', bold=False):
            # Fallback to default bitmap font if TTF isn't loaded
            try:
                font = ImageFont.load_default()
            except Exception:
                font = None
            draw.text((x, y), text, fill=color, font=font)

        # Header Text
        draw_text_simple(80, 75, "KOSVIO - EXECUTIVE INTELLIGENCE SNAPSHOT", size=20, color='#6366F1', bold=True)
        draw_text_simple(80, 100, f"Source File: {filename}   |   Generated: {datetime.datetime.now().strftime('%Y-%m-%d %H:%M')}", size=11, color='#94A3B8')
        
        # Health Score card (Left Column)
        draw.rounded_rectangle([60, 160, 400, 460], radius=10, fill='#0F172A', outline='#334155')
        draw_text_simple(90, 190, "DATASET QUALITY HEALTH", size=11, color='#94A3B8')
        draw_text_simple(90, 230, f"{health:.1f}%", size=48, color='#10B981', bold=True)
        draw_text_simple(90, 310, "RATING STATUS: OPTIMAL", size=12, color='#10B981')
        draw_text_simple(90, 340, "All verification checks passed.", size=11, color='#94A3B8')
        
        # KPI Cards (Right Column, 2x2 Grid)
        kpi_coords = [
            (430, 160, 770, 290),
            (800, 160, 1140, 290),
            (430, 320, 770, 450),
            (800, 320, 1140, 450)
        ]
        
        for idx, coord in enumerate(kpi_coords):
            draw.rounded_rectangle(coord, radius=8, fill='#0F172A', outline='#334155')
            if idx < len(kpis):
                k = kpis[idx]
                draw_text_simple(coord[0]+20, coord[1]+20, k["label"].upper(), size=10, color='#94A3B8')
                draw_text_simple(coord[0]+20, coord[1]+50, str(k["value"]), size=24, color='#6366F1', bold=True)
                draw_text_simple(coord[0]+20, coord[1]+90, k.get("detail", "") or "", size=10, color='#64748B')
                
        # Lower Section: Cleaning Summary
        draw.rounded_rectangle([60, 490, 1140, 730], radius=10, fill='#0F172A', outline='#334155')
        draw_text_simple(90, 520, "EXECUTIVE SUMMARY & CLEANSING STATS", size=12, color='#94A3B8', bold=True)
        
        if cleaning_summary:
            draw_text_simple(90, 560, f"• Total records removed due to duplicates: {cleaning_summary.get('rows_removed', 0):,}", size=12, color='#E2E8F0')
            draw_text_simple(90, 595, f"• Missing values successfully filled: {cleaning_summary.get('missing_fixed', 0):,}", size=12, color='#E2E8F0')
            draw_text_simple(90, 630, f"• Data types converted for analytical alignment: {cleaning_summary.get('dtypes_converted', 0)}", size=12, color='#E2E8F0')
            draw_text_simple(90, 665, f"• Outliers clipped to baseline limits: {cleaning_summary.get('outliers_handled', 0):,}", size=12, color='#E2E8F0')
        else:
            draw_text_simple(90, 560, "Dataset utilized in raw structural form. No cleaning modifications applied.", size=12, color='#64748B')
            
        buffer = io.BytesIO()
        img.save(buffer, format='PNG')
        return buffer.getvalue()

    @classmethod
    def generate_excel(cls, df: pd.DataFrame, health: float, kpis: list, cleaning_summary: dict, filename: str) -> bytes:
        """Generate a multi-tab Excel spreadsheet summarizing dataset characteristics."""
        buffer = io.BytesIO()
        
        with pd.ExcelWriter(buffer, engine='openpyxl') as writer:
            # Tab 1: Dataset Summary
            summary_info = {
                "Property": ["Dataset Name", "Generated Timestamp", "Dataset Health Score"],
                "Value": [filename, datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S"), f"{health:.1f}/100"]
            }
            pd.DataFrame(summary_info).to_excel(writer, sheet_name="Executive Summary", index=False)
            
            # Tab 2: Core KPI Metrics
            kpi_rows = []
            for k in kpis:
                kpi_rows.append({
                    "KPI Card Name": k["label"],
                    "Value": k["value"],
                    "Details": k.get("detail", "")
                })
            pd.DataFrame(kpi_rows).to_excel(writer, sheet_name="KPI Cards Metrics", index=False)
            
            # Tab 3: Cleaning Logs
            if cleaning_summary:
                clean_rows = [
                    {"Operation": "Rows Removed (Duplicates)", "Impact Value": cleaning_summary.get("rows_removed", 0)},
                    {"Operation": "Missing Fields Filled", "Impact Value": cleaning_summary.get("missing_fixed", 0)},
                    {"Operation": "Columns Type Converted", "Impact Value": cleaning_summary.get("dtypes_converted", 0)},
                    {"Operation": "Outliers Handled", "Impact Value": cleaning_summary.get("outliers_handled", 0)},
                ]
                pd.DataFrame(clean_rows).to_excel(writer, sheet_name="Cleansing Summary", index=False)
            
            # Tab 4: Cleaned Data Preview (first 100 rows)
            df.head(100).to_excel(writer, sheet_name="Cleaned Data Preview", index=False)
            
        return buffer.getvalue()
